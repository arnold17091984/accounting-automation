"""
P&L PowerPoint Builder Module

Generates formatted PowerPoint presentations for financial reports.
"""

import json
from datetime import datetime
from decimal import Decimal
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


class PLPowerPointBuilder:
    """Builds formatted PowerPoint presentations for P&L reports."""

    # Color scheme
    COLORS = {
        "primary": RGBColor(31, 78, 121),      # Dark blue
        "secondary": RGBColor(68, 114, 196),   # Medium blue
        "accent": RGBColor(255, 192, 0),       # Gold
        "positive": RGBColor(0, 176, 80),      # Green
        "negative": RGBColor(255, 0, 0),       # Red
        "neutral": RGBColor(128, 128, 128),    # Gray
        "white": RGBColor(255, 255, 255),
        "light_gray": RGBColor(242, 242, 242),
    }

    def __init__(self, config_dir: Path | str = None, template_dir: Path | str = None):
        """Initialize the PowerPoint builder.

        Args:
            config_dir: Path to configuration directory
            template_dir: Path to template directory
        """
        self.config_dir = Path(config_dir) if config_dir else Path(__file__).parent.parent.parent / "config"
        self.template_dir = Path(template_dir) if template_dir else Path(__file__).parent / "templates"
        self._load_config()

    def _load_config(self) -> None:
        """Load configuration files."""
        with open(self.config_dir / "entity_config.yaml") as f:
            self.entity_config = yaml.safe_load(f)

    def generate_monthly_report(
        self,
        entity: str,
        period: str,
        data: dict[str, Any],
        output_path: Path | str,
        template_path: Path | str = None,
        include_charts: bool = True,
        include_variance: bool = True,
    ) -> Path:
        """Generate a monthly P&L PowerPoint presentation.

        Args:
            entity: Entity code
            period: Period string (YYYY-MM)
            data: Financial data dictionary
            output_path: Output file path
            template_path: Optional template file path
            include_charts: Include chart slides
            include_variance: Include variance analysis

        Returns:
            Path to generated file
        """
        output_path = Path(output_path)

        # Load template or create new presentation
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
            prs.slide_height = Inches(7.5)

        # Get entity info
        entity_info = self.entity_config["entities"].get(entity, {})
        entity_name = entity_info.get("full_name", entity.title())
        period_date = datetime.strptime(period, "%Y-%m")

        # Build slides
        self._add_title_slide(prs, entity_name, period_date)
        self._add_executive_summary_slide(prs, data)
        self._add_pl_summary_slide(prs, data, entity_name)

        if include_charts:
            self._add_revenue_chart_slide(prs, data)
            self._add_expense_chart_slide(prs, data)

        if include_variance:
            self._add_variance_slide(prs, data)

        self._add_key_metrics_slide(prs, data)
        self._add_closing_slide(prs)

        # Save
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        return output_path

    def generate_consolidated_report(
        self,
        entities: list[str],
        period: str,
        data: dict[str, dict],
        output_path: Path | str,
    ) -> Path:
        """Generate a consolidated report for multiple entities.

        Args:
            entities: List of entity codes
            period: Period string (YYYY-MM)
            data: Dictionary of entity data {entity: financial_data}
            output_path: Output file path

        Returns:
            Path to generated file
        """
        output_path = Path(output_path)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        period_date = datetime.strptime(period, "%Y-%m")

        # Title slide
        self._add_consolidated_title_slide(prs, entities, period_date)

        # Consolidated summary
        self._add_consolidated_summary_slide(prs, data, entities)

        # Entity comparison
        self._add_entity_comparison_slide(prs, data, entities)

        # Individual entity summaries
        for entity in entities:
            entity_data = data.get(entity, {})
            if entity_data:
                entity_info = self.entity_config["entities"].get(entity, {})
                entity_name = entity_info.get("full_name", entity.title())
                self._add_entity_summary_slide(prs, entity_name, entity_data)

        self._add_closing_slide(prs)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        return output_path

    def _add_title_slide(self, prs: Presentation, entity_name: str, period_date: datetime) -> None:
        """Add the title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["primary"]
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = entity_name
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS["white"]
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = "Monthly Financial Report"
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = self.COLORS["accent"]
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Period
        period_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(12.333), Inches(0.75)
        )
        period_frame = period_box.text_frame
        period_para = period_frame.paragraphs[0]
        period_para.text = period_date.strftime("%B %Y")
        period_para.font.size = Pt(24)
        period_para.font.color.rgb = self.COLORS["white"]
        period_para.alignment = PP_ALIGN.CENTER

    def _add_executive_summary_slide(self, prs: Presentation, data: dict) -> None:
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        self._add_slide_title(slide, "Executive Summary")

        # Calculate key metrics
        total_revenue = sum(
            Decimal(str(item.get("actual", 0)))
            for item in data.get("revenue", {}).get("items", [])
        )
        total_cos = sum(
            Decimal(str(item.get("actual", 0)))
            for item in data.get("cost_of_sales", {}).get("items", [])
        )
        total_opex = sum(
            Decimal(str(item.get("actual", 0)))
            for item in data.get("operating_expenses", {}).get("items", [])
        )
        gross_profit = total_revenue - total_cos
        net_income = gross_profit - total_opex

        # Key metrics boxes
        metrics = [
            ("Revenue", total_revenue, data.get("revenue_change", 0)),
            ("Gross Profit", gross_profit, data.get("gross_profit_change", 0)),
            ("Net Income", net_income, data.get("net_income_change", 0)),
        ]

        box_width = Inches(3.5)
        box_height = Inches(2)
        start_x = Inches(1)
        start_y = Inches(2)
        gap = Inches(0.5)

        for i, (label, value, change) in enumerate(metrics):
            x = start_x + (box_width + gap) * i
            self._add_metric_box(slide, x, start_y, box_width, box_height, label, value, change)

        # Highlights section
        highlights = data.get("highlights", [])
        if highlights:
            highlights_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(11), Inches(2.5)
            )
            tf = highlights_box.text_frame
            tf.word_wrap = True

            title_para = tf.paragraphs[0]
            title_para.text = "Key Highlights"
            title_para.font.size = Pt(18)
            title_para.font.bold = True

            for highlight in highlights[:5]:
                para = tf.add_paragraph()
                para.text = f"• {highlight}"
                para.font.size = Pt(14)
                para.space_before = Pt(6)

    def _add_pl_summary_slide(self, prs: Presentation, data: dict, entity_name: str) -> None:
        """Add P&L summary table slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Profit & Loss Summary")

        # Create table
        rows = 12
        cols = 4
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5.5)
        ).table

        # Headers
        headers = ["Category", "Actual", "Budget", "Variance"]
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS["primary"]
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = self.COLORS["white"]
            para.font.bold = True
            para.font.size = Pt(12)

        # Data rows
        row_data = [
            ("Revenue", "revenue"),
            ("Cost of Sales", "cost_of_sales"),
            ("Gross Profit", None),
            ("Operating Expenses", "operating_expenses"),
            ("Net Income", None),
        ]

        current_row = 1
        for label, section_key in row_data:
            if section_key:
                items = data.get(section_key, {}).get("items", [])
                actual = sum(Decimal(str(i.get("actual", 0))) for i in items)
                budget = sum(Decimal(str(i.get("budget", 0))) for i in items)
            else:
                # Calculate totals
                if label == "Gross Profit":
                    rev = sum(Decimal(str(i.get("actual", 0))) for i in data.get("revenue", {}).get("items", []))
                    cos = sum(Decimal(str(i.get("actual", 0))) for i in data.get("cost_of_sales", {}).get("items", []))
                    actual = rev - cos
                    rev_b = sum(Decimal(str(i.get("budget", 0))) for i in data.get("revenue", {}).get("items", []))
                    cos_b = sum(Decimal(str(i.get("budget", 0))) for i in data.get("cost_of_sales", {}).get("items", []))
                    budget = rev_b - cos_b
                else:  # Net Income
                    rev = sum(Decimal(str(i.get("actual", 0))) for i in data.get("revenue", {}).get("items", []))
                    cos = sum(Decimal(str(i.get("actual", 0))) for i in data.get("cost_of_sales", {}).get("items", []))
                    opex = sum(Decimal(str(i.get("actual", 0))) for i in data.get("operating_expenses", {}).get("items", []))
                    actual = rev - cos - opex
                    rev_b = sum(Decimal(str(i.get("budget", 0))) for i in data.get("revenue", {}).get("items", []))
                    cos_b = sum(Decimal(str(i.get("budget", 0))) for i in data.get("cost_of_sales", {}).get("items", []))
                    opex_b = sum(Decimal(str(i.get("budget", 0))) for i in data.get("operating_expenses", {}).get("items", []))
                    budget = rev_b - cos_b - opex_b

            variance = actual - budget

            table.cell(current_row, 0).text = label
            table.cell(current_row, 1).text = f"₱{float(actual):,.0f}"
            table.cell(current_row, 2).text = f"₱{float(budget):,.0f}"
            table.cell(current_row, 3).text = f"₱{float(variance):,.0f}"

            # Highlight totals
            if label in ["Gross Profit", "Net Income"]:
                for col in range(4):
                    table.cell(current_row, col).fill.solid()
                    table.cell(current_row, col).fill.fore_color.rgb = self.COLORS["light_gray"]
                    para = table.cell(current_row, col).text_frame.paragraphs[0]
                    para.font.bold = True

            current_row += 1

    def _add_revenue_chart_slide(self, prs: Presentation, data: dict) -> None:
        """Add revenue breakdown chart slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Revenue Breakdown")

        # Revenue data
        revenue_items = data.get("revenue", {}).get("items", [])
        if not revenue_items:
            return

        # Create chart data
        chart_data = []
        for item in revenue_items:
            chart_data.append({
                "name": item.get("name", "Unknown"),
                "value": float(item.get("actual", 0))
            })

        # Add placeholder for chart (actual chart requires chart_data module)
        chart_placeholder = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11), Inches(4.5)
        )
        tf = chart_placeholder.text_frame
        para = tf.paragraphs[0]
        para.text = "Revenue Distribution"
        para.font.size = Pt(24)
        para.alignment = PP_ALIGN.CENTER

        # Add text representation of data
        for item in chart_data[:5]:
            para = tf.add_paragraph()
            para.text = f"• {item['name']}: ₱{item['value']:,.0f}"
            para.font.size = Pt(14)
            para.space_before = Pt(12)

    def _add_expense_chart_slide(self, prs: Presentation, data: dict) -> None:
        """Add expense breakdown chart slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Expense Breakdown")

        # Combine COS and OpEx
        all_expenses = (
            data.get("cost_of_sales", {}).get("items", []) +
            data.get("operating_expenses", {}).get("items", [])
        )

        if not all_expenses:
            return

        # Sort by amount
        sorted_expenses = sorted(
            all_expenses,
            key=lambda x: float(x.get("actual", 0)),
            reverse=True
        )

        # Add text representation
        expense_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11), Inches(5)
        )
        tf = expense_box.text_frame

        para = tf.paragraphs[0]
        para.text = "Top Expense Categories"
        para.font.size = Pt(20)
        para.font.bold = True

        for item in sorted_expenses[:8]:
            para = tf.add_paragraph()
            para.text = f"• {item.get('name', 'Unknown')}: ₱{float(item.get('actual', 0)):,.0f}"
            para.font.size = Pt(14)
            para.space_before = Pt(8)

    def _add_variance_slide(self, prs: Presentation, data: dict) -> None:
        """Add budget variance analysis slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Budget Variance Analysis")

        # Find significant variances
        all_items = (
            data.get("revenue", {}).get("items", []) +
            data.get("cost_of_sales", {}).get("items", []) +
            data.get("operating_expenses", {}).get("items", [])
        )

        variances = []
        for item in all_items:
            actual = float(item.get("actual", 0))
            budget = float(item.get("budget", 0))
            if budget > 0:
                variance_pct = ((actual - budget) / budget) * 100
                if abs(variance_pct) > 10:
                    variances.append({
                        "name": item.get("name", "Unknown"),
                        "actual": actual,
                        "budget": budget,
                        "variance_pct": variance_pct
                    })

        variances.sort(key=lambda x: abs(x["variance_pct"]), reverse=True)

        # Add variance list
        variance_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(12), Inches(5)
        )
        tf = variance_box.text_frame

        if variances:
            para = tf.paragraphs[0]
            para.text = "Significant Variances (>10%)"
            para.font.size = Pt(16)
            para.font.bold = True

            for item in variances[:8]:
                para = tf.add_paragraph()
                sign = "+" if item["variance_pct"] > 0 else ""
                status = "⚠️ Over" if item["variance_pct"] > 0 else "✓ Under"
                para.text = f"{status} {item['name']}: {sign}{item['variance_pct']:.1f}% (₱{item['actual']:,.0f} vs ₱{item['budget']:,.0f})"
                para.font.size = Pt(12)
                para.space_before = Pt(8)

                if item["variance_pct"] > 0:
                    para.font.color.rgb = self.COLORS["negative"]
                else:
                    para.font.color.rgb = self.COLORS["positive"]
        else:
            para = tf.paragraphs[0]
            para.text = "✓ All categories within 10% of budget"
            para.font.size = Pt(18)
            para.font.color.rgb = self.COLORS["positive"]

    def _add_key_metrics_slide(self, prs: Presentation, data: dict) -> None:
        """Add key metrics slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Key Performance Indicators")

        # Calculate metrics
        total_revenue = sum(
            Decimal(str(item.get("actual", 0)))
            for item in data.get("revenue", {}).get("items", [])
        )
        total_cos = sum(
            Decimal(str(item.get("actual", 0)))
            for item in data.get("cost_of_sales", {}).get("items", [])
        )
        total_opex = sum(
            Decimal(str(item.get("actual", 0)))
            for item in data.get("operating_expenses", {}).get("items", [])
        )
        gross_profit = total_revenue - total_cos
        net_income = gross_profit - total_opex

        gross_margin = (gross_profit / total_revenue * 100) if total_revenue else 0
        net_margin = (net_income / total_revenue * 100) if total_revenue else 0
        opex_ratio = (total_opex / total_revenue * 100) if total_revenue else 0

        metrics = [
            ("Gross Margin", f"{float(gross_margin):.1f}%"),
            ("Net Margin", f"{float(net_margin):.1f}%"),
            ("OpEx Ratio", f"{float(opex_ratio):.1f}%"),
        ]

        # Add metric boxes
        box_width = Inches(3.5)
        box_height = Inches(1.5)
        start_x = Inches(1)
        start_y = Inches(2.5)
        gap = Inches(0.5)

        for i, (label, value) in enumerate(metrics):
            x = start_x + (box_width + gap) * i

            box = slide.shapes.add_shape(
                1, x, start_y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS["secondary"]

            # Label
            label_box = slide.shapes.add_textbox(
                x, start_y + Inches(0.2), box_width, Inches(0.5)
            )
            label_tf = label_box.text_frame
            label_para = label_tf.paragraphs[0]
            label_para.text = label
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = self.COLORS["white"]
            label_para.alignment = PP_ALIGN.CENTER

            # Value
            value_box = slide.shapes.add_textbox(
                x, start_y + Inches(0.6), box_width, Inches(0.7)
            )
            value_tf = value_box.text_frame
            value_para = value_tf.paragraphs[0]
            value_para.text = value
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            value_para.font.color.rgb = self.COLORS["white"]
            value_para.alignment = PP_ALIGN.CENTER

    def _add_closing_slide(self, prs: Presentation) -> None:
        """Add closing slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["primary"]
        background.line.fill.background()

        # Thank you text
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        tf = text_box.text_frame
        para = tf.paragraphs[0]
        para.text = "Thank You"
        para.font.size = Pt(48)
        para.font.bold = True
        para.font.color.rgb = self.COLORS["white"]
        para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.75)
        )
        sub_tf = subtitle_box.text_frame
        sub_para = sub_tf.paragraphs[0]
        sub_para.text = "Generated by BK Keyforce Accounting Automation"
        sub_para.font.size = Pt(16)
        sub_para.font.color.rgb = self.COLORS["accent"]
        sub_para.alignment = PP_ALIGN.CENTER

    def _add_slide_title(self, slide, title: str) -> None:
        """Add a slide title."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.text = title
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = self.COLORS["primary"]

    def _add_metric_box(
        self,
        slide,
        x: Inches,
        y: Inches,
        width: Inches,
        height: Inches,
        label: str,
        value: Decimal,
        change: float = 0
    ) -> None:
        """Add a metric display box."""
        # Box background
        box = slide.shapes.add_shape(1, x, y, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS["light_gray"]

        # Label
        label_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(0.4))
        label_tf = label_box.text_frame
        label_para = label_tf.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = self.COLORS["neutral"]
        label_para.alignment = PP_ALIGN.CENTER

        # Value
        value_box = slide.shapes.add_textbox(x, y + Inches(0.6), width, Inches(0.8))
        value_tf = value_box.text_frame
        value_para = value_tf.paragraphs[0]
        value_para.text = f"₱{float(value):,.0f}"
        value_para.font.size = Pt(24)
        value_para.font.bold = True
        value_para.font.color.rgb = self.COLORS["primary"]
        value_para.alignment = PP_ALIGN.CENTER

        # Change indicator
        if change != 0:
            change_box = slide.shapes.add_textbox(x, y + Inches(1.4), width, Inches(0.4))
            change_tf = change_box.text_frame
            change_para = change_tf.paragraphs[0]
            sign = "+" if change > 0 else ""
            change_para.text = f"{sign}{change:.1f}%"
            change_para.font.size = Pt(12)
            change_para.font.color.rgb = self.COLORS["positive"] if change >= 0 else self.COLORS["negative"]
            change_para.alignment = PP_ALIGN.CENTER

    def _add_consolidated_title_slide(self, prs: Presentation, entities: list[str], period_date: datetime) -> None:
        """Add consolidated report title slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["primary"]
        background.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.text = "Consolidated Financial Report"
        para.font.size = Pt(48)
        para.font.bold = True
        para.font.color.rgb = self.COLORS["white"]
        para.alignment = PP_ALIGN.CENTER

        entities_text = ", ".join([e.title() for e in entities])
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12.333), Inches(1)
        )
        sub_tf = subtitle_box.text_frame
        sub_para = sub_tf.paragraphs[0]
        sub_para.text = entities_text
        sub_para.font.size = Pt(20)
        sub_para.font.color.rgb = self.COLORS["accent"]
        sub_para.alignment = PP_ALIGN.CENTER

        period_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(12.333), Inches(0.75)
        )
        period_tf = period_box.text_frame
        period_para = period_tf.paragraphs[0]
        period_para.text = period_date.strftime("%B %Y")
        period_para.font.size = Pt(18)
        period_para.font.color.rgb = self.COLORS["white"]
        period_para.alignment = PP_ALIGN.CENTER

    def _add_consolidated_summary_slide(self, prs: Presentation, data: dict, entities: list[str]) -> None:
        """Add consolidated summary slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Consolidated Summary")

        # Calculate totals across all entities
        total_revenue = Decimal("0")
        total_net_income = Decimal("0")

        for entity in entities:
            entity_data = data.get(entity, {})
            rev = sum(Decimal(str(i.get("actual", 0))) for i in entity_data.get("revenue", {}).get("items", []))
            cos = sum(Decimal(str(i.get("actual", 0))) for i in entity_data.get("cost_of_sales", {}).get("items", []))
            opex = sum(Decimal(str(i.get("actual", 0))) for i in entity_data.get("operating_expenses", {}).get("items", []))
            total_revenue += rev
            total_net_income += (rev - cos - opex)

        # Add summary boxes
        metrics = [
            ("Total Revenue", total_revenue),
            ("Total Net Income", total_net_income),
            ("Entities", len(entities)),
        ]

        box_width = Inches(3.5)
        start_x = Inches(1)
        start_y = Inches(2.5)
        gap = Inches(0.5)

        for i, (label, value) in enumerate(metrics):
            x = start_x + (box_width + gap) * i
            if isinstance(value, int):
                self._add_simple_metric_box(slide, x, start_y, box_width, Inches(1.5), label, str(value))
            else:
                self._add_metric_box(slide, x, start_y, box_width, Inches(1.5), label, value)

    def _add_simple_metric_box(self, slide, x, y, width, height, label: str, value: str) -> None:
        """Add a simple metric box with string value."""
        box = slide.shapes.add_shape(1, x, y, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS["light_gray"]

        label_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(0.4))
        label_tf = label_box.text_frame
        label_para = label_tf.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(14)
        label_para.alignment = PP_ALIGN.CENTER

        value_box = slide.shapes.add_textbox(x, y + Inches(0.6), width, Inches(0.6))
        value_tf = value_box.text_frame
        value_para = value_tf.paragraphs[0]
        value_para.text = value
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.alignment = PP_ALIGN.CENTER

    def _add_entity_comparison_slide(self, prs: Presentation, data: dict, entities: list[str]) -> None:
        """Add entity comparison slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Entity Comparison")

        # Comparison text box
        comparison_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(12), Inches(5)
        )
        tf = comparison_box.text_frame

        for entity in entities:
            entity_info = self.entity_config["entities"].get(entity, {})
            entity_name = entity_info.get("full_name", entity.title())
            entity_data = data.get(entity, {})

            rev = sum(Decimal(str(i.get("actual", 0))) for i in entity_data.get("revenue", {}).get("items", []))
            cos = sum(Decimal(str(i.get("actual", 0))) for i in entity_data.get("cost_of_sales", {}).get("items", []))
            opex = sum(Decimal(str(i.get("actual", 0))) for i in entity_data.get("operating_expenses", {}).get("items", []))
            net = rev - cos - opex

            para = tf.add_paragraph()
            para.text = f"{entity_name}: Revenue ₱{float(rev):,.0f} | Net Income ₱{float(net):,.0f}"
            para.font.size = Pt(14)
            para.space_before = Pt(12)

    def _add_entity_summary_slide(self, prs: Presentation, entity_name: str, data: dict) -> None:
        """Add individual entity summary slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, entity_name)

        rev = sum(Decimal(str(i.get("actual", 0))) for i in data.get("revenue", {}).get("items", []))
        cos = sum(Decimal(str(i.get("actual", 0))) for i in data.get("cost_of_sales", {}).get("items", []))
        opex = sum(Decimal(str(i.get("actual", 0))) for i in data.get("operating_expenses", {}).get("items", []))
        gross = rev - cos
        net = gross - opex

        summary_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11), Inches(4)
        )
        tf = summary_box.text_frame

        metrics = [
            ("Revenue", rev),
            ("Cost of Sales", cos),
            ("Gross Profit", gross),
            ("Operating Expenses", opex),
            ("Net Income", net),
        ]

        for label, value in metrics:
            para = tf.add_paragraph()
            para.text = f"{label}: ₱{float(value):,.0f}"
            para.font.size = Pt(18)
            para.space_before = Pt(12)

            if label == "Net Income":
                para.font.bold = True
                para.font.color.rgb = self.COLORS["positive"] if value >= 0 else self.COLORS["negative"]


def main():
    """CLI entry point for testing."""
    import argparse

    parser = argparse.ArgumentParser(description="Generate P&L PowerPoint")
    parser.add_argument("--entity", required=True, help="Entity code")
    parser.add_argument("--month", required=True, help="Period (YYYY-MM)")
    parser.add_argument("--data-file", required=True, help="JSON data file")
    parser.add_argument("--output-dir", default="/tmp/output", help="Output directory")

    args = parser.parse_args()

    with open(args.data_file) as f:
        data = json.load(f)

    builder = PLPowerPointBuilder()
    output_path = Path(args.output_dir) / f"{args.entity}_pl_{args.month}.pptx"
    result = builder.generate_monthly_report(
        entity=args.entity,
        period=args.month,
        data=data,
        output_path=output_path
    )

    print(f"Generated: {result}")


if __name__ == "__main__":
    main()
